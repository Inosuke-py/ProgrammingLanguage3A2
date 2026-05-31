"""PPTX parsing service using python-pptx."""

from pptx import Presentation


def parse_pptx(file_path: str) -> list[dict]:
    """
    Parse a PPTX file into sections.

    Each section represents one slide. Returns a list of dicts with:
    - title: slide title (from the title placeholder) or None
    - content: all text extracted from the slide's shapes
    - page_number: slide number (1-indexed)
    """
    prs = Presentation(file_path)
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_title = None
        text_parts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Check if this shape is the title placeholder
            if shape.shape_type == 13 or (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            ):
                title_text = shape.text_frame.text.strip()
                if title_text:
                    slide_title = title_text
                    continue

            # Extract all text from the shape's paragraphs
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    text_parts.append(para_text)

        content = "\n".join(text_parts).strip()

        # Skip slides with no text content
        if not content and not slide_title:
            continue

        sections.append({
            "title": slide_title,
            "content": content if content else (slide_title or ""),
            "page_number": slide_num,
        })

    return sections
